import io
import PyPDF2
import docx
try:
    from pptx import Presentation as PptxPresentation
    PPTX_SUPPORTED = True
except ImportError:
    PPTX_SUPPORTED = False

def extract_text_from_file(file_bytes: bytes, filename: str) -> tuple[str, str]:
    name_lower = filename.lower()
    if name_lower.endswith(".pdf"):
        return _extract_pdf(file_bytes), "pdf"
    elif name_lower.endswith(".docx"):
        return _extract_docx(file_bytes), "docx"
    elif name_lower.endswith(".pptx") and PPTX_SUPPORTED:
        return _extract_pptx(file_bytes), "pptx"
    elif name_lower.endswith(".txt") or name_lower.endswith(".md"):
        return file_bytes.decode("utf-8", errors="ignore"), "txt"
    else:
        try:
            return file_bytes.decode("utf-8", errors="ignore"), "txt"
        except:
            raise ValueError(f"Unsupported file type: {filename}")

def _extract_pdf(file_bytes: bytes) -> str:
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text)
    return "\n".join(texts)

def _extract_docx(file_bytes: bytes) -> str:
    doc = docx.Document(io.BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

def _extract_pptx(file_bytes: bytes) -> str:
    prs = PptxPresentation(io.BytesIO(file_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)